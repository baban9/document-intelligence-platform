"""Text extraction handlers per document kind."""

from __future__ import annotations

import csv
import json
from io import StringIO
from pathlib import Path

import fitz

from docintel.capabilities.extraction.formats.models import DocumentKind, ExtractionResult, IdentificationResult
from docintel.capabilities.extraction.formats.sniff import identify_document


def extract_document_text(
    path: str | Path,
    *,
    filename: str | None = None,
    content_type: str | None = None,
    identification: IdentificationResult | None = None,
) -> ExtractionResult:
    """Extract plain text from a supported document upload."""
    file_path = Path(path)
    resolved = identification or identify_document(
        file_path,
        filename=filename,
        content_type=content_type,
    )
    if resolved.profile is None or not resolved.profile.supports_text_extraction:
        raise ValueError(
            f"Unsupported document kind '{resolved.kind.value}'. "
            "Call GET /v1/documents/types for supported formats."
        )

    if resolved.kind is DocumentKind.PDF:
        return _extract_pdf(file_path, resolved)
    if resolved.kind is DocumentKind.DOCX:
        return _extract_docx(file_path, resolved)
    if resolved.kind is DocumentKind.XLSX:
        return _extract_xlsx(file_path, resolved)
    if resolved.kind is DocumentKind.PPTX:
        return _extract_pptx(file_path, resolved)
    if resolved.kind is DocumentKind.CSV:
        return _extract_csv(file_path, resolved)
    if resolved.kind is DocumentKind.JSON:
        return _extract_json(file_path, resolved)
    if resolved.kind is DocumentKind.PLAIN_TEXT:
        return _extract_plain_text(file_path, resolved)

    raise ValueError(f"No extractor registered for kind '{resolved.kind.value}'.")


def _extract_pdf(path: Path, identification: IdentificationResult) -> ExtractionResult:
    from docintel.capabilities.extraction.formats.paginated_pdf import extract_pdf_document

    return extract_pdf_document(path, identification)


def _extract_docx(path: Path, identification: IdentificationResult) -> ExtractionResult:
    try:
        import docx
    except ImportError as exc:
        raise RuntimeError(
            "DOCX support requires optional dependencies. Install: pip install -e '.[documents]'"
        ) from exc

    document = docx.Document(path)
    parts: list[str] = []
    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return ExtractionResult(
        kind=identification.kind,
        mime_type=identification.mime_type,
        text="\n".join(parts),
        segments=[{"section": "body", "text": "\n".join(parts)}],
        metadata={"paragraph_count": len(document.paragraphs), "table_count": len(document.tables)},
    )


def _extract_xlsx(path: Path, identification: IdentificationResult) -> ExtractionResult:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise RuntimeError(
            "Excel support requires optional dependencies. Install: pip install -e '.[documents]'"
        ) from exc

    workbook = load_workbook(path, read_only=True, data_only=True)
    segments: list[dict] = []
    parts: list[str] = []
    for sheet in workbook.worksheets:
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(value).strip() for value in row if value is not None and str(value).strip()]
            if cells:
                rows.append(" | ".join(cells))
        sheet_text = "\n".join(rows)
        segments.append({"sheet": sheet.title, "text": sheet_text})
        if sheet_text:
            parts.append(f"# {sheet.title}\n{sheet_text}")
    workbook.close()

    return ExtractionResult(
        kind=identification.kind,
        mime_type=identification.mime_type,
        text="\n\n".join(parts),
        segments=segments,
        metadata={"sheet_count": len(segments)},
    )


def _extract_pptx(path: Path, identification: IdentificationResult) -> ExtractionResult:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError(
            "PowerPoint support requires optional dependencies. Install: pip install -e '.[documents]'"
        ) from exc

    presentation = Presentation(path)
    segments: list[dict] = []
    parts: list[str] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                slide_parts.append(text)
        slide_text = "\n".join(slide_parts)
        segments.append({"slide": slide_index, "text": slide_text})
        if slide_text:
            parts.append(f"# Slide {slide_index}\n{slide_text}")

    return ExtractionResult(
        kind=identification.kind,
        mime_type=identification.mime_type,
        text="\n\n".join(parts),
        segments=segments,
        metadata={"slide_count": len(presentation.slides)},
    )


def _extract_csv(path: Path, identification: IdentificationResult) -> ExtractionResult:
    raw = path.read_text(encoding="utf-8", errors="replace")
    sample = raw[:2048]
    try:
        dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
    except csv.Error:
        dialect = csv.excel

    reader = csv.reader(StringIO(raw), dialect)
    rows = [" | ".join(cell.strip() for cell in row if cell.strip()) for row in reader]
    rows = [row for row in rows if row]
    return ExtractionResult(
        kind=identification.kind,
        mime_type=identification.mime_type,
        text="\n".join(rows),
        segments=[{"section": "rows", "text": "\n".join(rows)}],
        metadata={"row_count": len(rows)},
    )


def _extract_json(path: Path, identification: IdentificationResult) -> ExtractionResult:
    payload = json.loads(path.read_text(encoding="utf-8"))
    return ExtractionResult(
        kind=identification.kind,
        mime_type=identification.mime_type,
        text=json.dumps(payload, indent=2, sort_keys=True),
        segments=[{"section": "json", "text": json.dumps(payload)}],
        metadata={"top_level_type": type(payload).__name__},
    )


def _extract_plain_text(path: Path, identification: IdentificationResult) -> ExtractionResult:
    text = path.read_text(encoding="utf-8", errors="replace")
    return ExtractionResult(
        kind=identification.kind,
        mime_type=identification.mime_type,
        text=text,
        segments=[{"section": "body", "text": text}],
        metadata={"char_count": len(text)},
    )
